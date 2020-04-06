import requests
from bs4 import BeautifulSoup
from googlesearch import search 
from flask import Flask
from flask import request as flrequest
from selenium import webdriver 
from selenium.webdriver.common.by import By 
from selenium.webdriver.support.ui import WebDriverWait 
from selenium.webdriver.support import expected_conditions as EC 
from selenium.common.exceptions import TimeoutException
from flask import jsonify
import pyrebase
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
import codecs
import cv2
import numpy as np
import argparse
from sklearn.cluster import KMeans
from collections import Counter
import math
import os
import json
import time
import uuid
import os.path
from os import path






path_to_here = os.path.dirname(os.path.abspath(__file__))

if(not path.exists(os.path.join(path_to_here, "images"))):
    os.mkdir(os.path.join(path_to_here, "images"))




config = {
  "apiKey": "",
  "authDomain": "<App Project ID>.firebaseapp.com",
  "databaseURL": "",
  "storageBucket": ""
}

firebase = pyrebase.initialize_app(config)

fbdb = firebase.database()
fbst = firebase.storage()
auth = firebase.auth()
user = auth.sign_in_with_email_and_password("server@server.com", "server123@12")

IP_Address = fbdb.get()

# Complementary Colours
def complement(r, g, b):
    #k = hilo(r, g, b)
    return tuple(abs(255 - u) for u in (r, g, b))


def Background_Image(slide, BackgroundPicturePath, prs):
    #change background with an image of the slide
    left = top = 0
    pic = slide.shapes.add_picture(BackgroundPicturePath, left, top, height = prs.slide_height)
    # move picture to background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)  # use the number that does the appropriate job
    return

def get_dominant_color(image, k=4, image_processing_size = None):
    """
    takes an image as input
    returns the dominant color of the image as a list
    
    dominant color is found by running k means on the 
    pixels & returning the centroid of the largest cluster

    processing time is sped up by working with a smaller image; 
    this resizing can be done with the image_processing_size param 
    which takes a tuple of image dims as input

    >>> get_dominant_color(my_image, k=4, image_processing_size = (25, 25))
    [56.2423442, 34.0834233, 70.1234123]
    """
    #resize image if new dims provided
    if image_processing_size is not None:
        image = cv2.resize(image, image_processing_size, 
                            interpolation = cv2.INTER_AREA)
    
    #reshape the image to be a list of pixels
    image = image.reshape((image.shape[0] * image.shape[1], 3))

    #cluster and assign labels to the pixels 
    clt = KMeans(n_clusters = k)
    labels = clt.fit_predict(image)

    #count labels to find most popular
    label_counts = Counter(labels)

    #subset out most popular centroid
    dominant_color = clt.cluster_centers_[label_counts.most_common(1)[0][0]]

    return list(dominant_color)


def Lyric_Resultify(query):
    azlyricsquery = search((query + "azlyrics"), tld="co.in", num=1, stop=1)
    azlyricslink = ""
    for j in azlyricsquery:
        azlyricslink = j
    URL = azlyricslink
    page = requests.get(URL)

    soup = BeautifulSoup(page.content, 'html.parser')

    inter = soup.find('div', class_="ringtone")
    final_lyrics = inter.find_next_sibling('div').text
    return final_lyrics

def Reading_Resultify(choice):
    path = r'C:\\Users\\Jerry Allan Akshay\\Desktop\\chromedriver_win32\\chromedriver.exe'



    browser = webdriver.Chrome(executable_path=path)


    Sunday_Reading_URL = "https://www.ballymenaparish.org/sundays-mass-readings/?feature=sunday" #Choice 1

    Today_Reading_URL = "https://www.ballymenaparish.org/sundays-mass-readings/?feature=today" #Choice 2


    if(choice == "sunday"):
        browser.get(Sunday_Reading_URL)
    else:
        browser.get(Today_Reading_URL)

    # Wait 30 seconds for page to load
    timeout = 30
    try:
        #WebDriverWait(browser, timeout).until(EC.visibility_of_element_located((By.XPATH, "//span[@style=\"font-family: Times New Roman;\"]")))
        WebDriverWait(browser, timeout).until(EC.visibility_of_element_located((By.XPATH, "//p")))
    except TimeoutException:
        print("Timed out waiting for page to load")
        browser.quit()

    readings = browser.find_element_by_id("ci-readings")
    final = readings.text.replace(u'\u2019', '\'').replace(u'\u2018', '\'').replace(u'\u201c', '\"').replace(u'\u201d', '\"').encode("utf-8")
    browser.close()

    splitstring = final.split()
    readings = []
    gospel = []
    response = []
    count = 0
    gospel_count = 0
    reading_count = 0
    response_count = 0
    read_start_count = read_end_count = 0
    gosp_start_count = gosp_end_count = 0
    resp_start_count = resp_end_count  = 0
    acc_start_count = acc_end_count  = 0
    cur_reading = False
    gosp_reading = False
    resp_reading = False
    acc_reading = False
    for word in splitstring:
        
        if (acc_reading==False and word.lower().strip()=="gospel" and splitstring[count+1].lower().strip() == "acclamation"):
            acc_reading = True
            acc_start_count = count+2

        if (cur_reading==False and word.lower().strip()=="reading"):
            cur_reading = True
            read_start_count = count+1

        if (resp_reading==False and word.lower().strip()=="responsorial"):
            resp_reading = True
            resp_start_count = count+2

        if (gosp_reading==False and word.lower().strip()=="gospel"):
            if(splitstring[count+1].lower().strip() == "acclamation"):
                count = count + 1
                continue
            elif(splitstring[count+1].lower().strip() == "reflection"):
                break
            gosp_reading = True
            gosp_start_count = count+1


        if(cur_reading==True and splitstring[count].lower().strip()=="the" and splitstring[count+1].lower().strip()=="word" and splitstring[count+2].lower().strip()=="of" and splitstring[count+3].lower().strip()=="the" and (splitstring[count+4].lower().strip()=="lord" or splitstring[count+4].lower().strip()=="lord.")):
            readings.append("")
            cur_reading = False
            read_end_count = count
            for i in range(read_start_count, read_end_count):
                readings[reading_count] = readings[reading_count] + splitstring[i] + " "
            readings[reading_count] = readings[reading_count] + "\n\n" + "The Word of the Lord."
            reading_count = reading_count + 1

        if(gosp_reading==True and splitstring[count].lower().strip()=="the" and splitstring[count+1].lower().strip()=="gospel" and splitstring[count+2].lower().strip()=="of" and splitstring[count+3].lower().strip()=="the" and (splitstring[count+4].lower().strip()=="lord" or splitstring[count+4].lower().strip()=="lord.")):
            gospel.append("")
            gosp_reading = False
            gosp_end_count = count
            for i in range(gosp_start_count, gosp_end_count):
                gospel[gospel_count] = gospel[gospel_count] + splitstring[i] + " "
            gospel[gospel_count] = gospel[gospel_count] +"\n\n" + "The Gospel of the Lord." 
            gospel_count = gospel_count + 1

        if(resp_reading==True and (splitstring[count].lower().strip()=="second" or splitstring[count].lower().strip()=="gospel" or splitstring[count].lower().strip()=="first" or splitstring[count].lower().strip()=="third" or splitstring[count].lower().strip()=="fourth" or splitstring[count].lower().strip()=="fifth")):
            response.append("")
            resp_reading = False
            resp_end_count = count
            for i in range(resp_start_count, resp_end_count):
                if(splitstring[i][0].isdigit()):
                    response[response_count] = response[response_count] + "\n\n"
                response[response_count] = response[response_count] + splitstring[i] + " "

            response_count = response_count + 1

        if(acc_reading==True and splitstring[count].lower().strip()=="gospel"):
            gospel.append("")
            gospel[gospel_count] = gospel[gospel_count] + "Gospel Acclamation (Keep this as the title)\n\n" 
            acc_reading = False
            acc_end_count = count
            for i in range(acc_start_count, acc_end_count):
                gospel[gospel_count] = gospel[gospel_count] + splitstring[i] + " "
            gospel[gospel_count] = gospel[gospel_count].strip() 
            gospel_count = gospel_count + 1
            
        count = count + 1
    
    return jsonify(readings=readings, gospels=gospel, response=response)


def Generate_PPT(ppt_datax):
    print(ppt_datax)
    ppt_data = json.loads(json.dumps(ppt_datax))
    #Global Parameters
    prs = Presentation()
    text_colours = []
    BackgroundPicturePath = "Black_Background_Default.jpg"
    BackgroundPicturePath = os.path.join(path_to_here, BackgroundPicturePath)
    if ppt_data['does_background_exist'] == True:
        fbst.child(ppt_data['background_image_path']).download(os.path.join(path_to_here, ppt_data['background_image_path'].replace('/', "\\")))
        BackgroundPicturePath = os.path.join(path_to_here, ppt_data['background_image_path'])
        bgimage = cv2.imread(BackgroundPicturePath)
        width = 960
        height = 720
        dim = (width, height)
        resized = cv2.resize(bgimage, dim, interpolation = cv2.INTER_CUBIC)
        cv2.imwrite(BackgroundPicturePath, resized)
    TitleFontColour = ContentFontColour = []
    if ppt_data['use_generated_colours'] == True:
        text_colours = get_dominant_color(cv2.cvtColor(cv2.imread(BackgroundPicturePath), cv2.COLOR_BGR2HSV))
        text_colours = complement(int(text_colours[0]), int(text_colours[1]), int(text_colours[2]))
        TitleFontColour = text_colours
        ContentFontColour = text_colours
    else:
        TitleFontColour = ppt_data['title_colour']
        ContentFontColour = ppt_data['content_colour']

    #Title Page
    MainTitlePageContent = ppt_data["main_title"]
    MainTitlePageSubtitle = ppt_data["main_subtitle"]


    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    Background_Image(slide, BackgroundPicturePath, prs)
    tf = title.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = MainTitlePageContent
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(60)
    font.bold = True
    font.color.rgb = RGBColor(int(TitleFontColour[0]), int(TitleFontColour[1]), int(TitleFontColour[2]))


    tf = subtitle.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = MainTitlePageSubtitle
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(50)
    font.bold = True
    run.font.color.rgb = RGBColor(int(ContentFontColour[0]), int(ContentFontColour[1]), int(ContentFontColour[2]))


    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    Background_Image(slide, BackgroundPicturePath, prs)
    title = slide.shapes.title
    body = slide.placeholders[1]

    #Saving the ppt
    file_path = str(uuid.uuid1()) + ".pptx" #(ppt_data["background_image_path"][7:ppt_data["background_image_path"].rindex('.')] + ".pptx").replace(':','_')
    print(file_path)
    prs.save(os.path.join(path_to_here, file_path))

    fbst.child("files/"+file_path).put(os.path.join(path_to_here, file_path), user['idToken'])

    os.remove(os.path.join(path_to_here, file_path))
    
    return ("files/"+file_path)



app = Flask(__name__)



@app.route("/")
def home():
    return "You've reached the server. Dance around now!"

@app.route("/lyrics")
def lyrics():
    lyquery = flrequest.args.get('lyquery')
    return jsonify(lyrics=Lyric_Resultify(lyquery))

@app.route("/readings")
def readings():
    rdchoice = flrequest.args.get('rdchoice')
    return Reading_Resultify(rdchoice)

@app.route("/createppt", methods=['GET', 'POST'])
def createppt():
    if(flrequest.method == 'POST'):
        ppt_data = flrequest.get_json()
        return jsonify(download_url=Generate_PPT(ppt_data))
    
if __name__ == "__main__":
    app.run(host = '192.168.1.126', port='5000')
    



