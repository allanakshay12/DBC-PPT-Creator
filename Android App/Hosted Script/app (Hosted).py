
import requests
from bs4 import BeautifulSoup
from googlesearch import search
from flask import Flask
from flask import request as flrequest
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.common.exceptions import TimeoutException
from flask import jsonify
import pyrebase
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
import codecs
import cv2
import numpy as np
import argparse
from sklearn.cluster import KMeans
from collections import Counter
import math
import os
import json
import time
import uuid
import os.path
from os import path
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
import copy 
import six

#GOOGLE_CHROME_PATH = '/app/.apt/usr/bin/google_chrome'
#CHROMEDRIVER_PATH = '/app/.chromedriver/bin/chromedriver'







path_to_here = os.path.dirname(os.path.abspath(__file__))

if(not path.exists(os.path.join(path_to_here, "images"))):
    os.mkdir(os.path.join(path_to_here, "images"))




config = {
  "apiKey": "",
  "authDomain": "<App Project ID>.firebaseapp.com",
  "databaseURL": "",
  "storageBucket": ""
}

firebase = pyrebase.initialize_app(config)

fbdb = firebase.database()
fbst = firebase.storage()
auth = firebase.auth()
user = auth.sign_in_with_email_and_password("server@server.com", "server123@12")


# Complementary Colours
def complement(r, g, b):
    #k = hilo(r, g, b)
    return tuple(abs(255 - u) for u in (r, g, b))


def Background_Image(slide, BackgroundPicturePath, prs):
    #change background with an image of the slide
    left = top = 0
    pic = slide.shapes.add_picture(BackgroundPicturePath, left, top, height = prs[0].slide_height)
    # move picture to background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)  # use the number that does the appropriate job
    return

def get_dominant_color(image, k=4, image_processing_size = None):
    """
    takes an image as input
    returns the dominant color of the image as a list

    dominant color is found by running k means on the
    pixels & returning the centroid of the largest cluster

    processing time is sped up by working with a smaller image;
    this resizing can be done with the image_processing_size param
    which takes a tuple of image dims as input

    >>> get_dominant_color(my_image, k=4, image_processing_size = (25, 25))
    [56.2423442, 34.0834233, 70.1234123]
    """
    #resize image if new dims provided
    if image_processing_size is not None:
        image = cv2.resize(image, image_processing_size,
                            interpolation = cv2.INTER_AREA)

    #reshape the image to be a list of pixels
    image = image.reshape((image.shape[0] * image.shape[1], 3))

    #cluster and assign labels to the pixels
    clt = KMeans(n_clusters = k)
    labels = clt.fit_predict(image)

    #count labels to find most popular
    label_counts = Counter(labels)

    #subset out most popular centroid
    dominant_color = clt.cluster_centers_[label_counts.most_common(1)[0][0]]

    return list(dominant_color)


def Lyric_Resultify(query):
    azlyricsquery = search(("site:songlyrics.com " + query), tld="com", num=10, stop=10)
    for j in azlyricsquery:
        try:
            URL = j
            page = requests.get(URL, timeout = 60000)
            print(URL)
            soup = BeautifulSoup(page.content, 'html.parser')
            print(soup)
            lyrics = soup.find('p', id="songLyricsDiv")
            return lyrics.text
            break
        except Exception:
            continue
    return "No lyrics found. Try a different title."

def Reading_Resultify(choice):

    Sunday_Reading_URL = "https://www.ballymenaparish.org/sundays-mass-readings/?feature=sunday" #Choice 1

    Today_Reading_URL = "https://www.ballymenaparish.org/sundays-mass-readings/?feature=today" #Choice 2


    final = ""

    chrome_options = webdriver.ChromeOptions()
    chrome_options.add_argument('--disable-dev-shm-usage')
    chrome_options.add_argument('--no-sandbox')
    chrome_options.add_argument('--headless')
    chrome_options.binary_location = os.environ.get("GOOGLE_CHROME_BIN")
    browser = webdriver.Chrome(executable_path=os.environ.get("CHROMEDRIVER_PATH"), chrome_options=chrome_options)
    browser.set_page_load_timeout(60000)
    if(choice == "sunday"):
        browser.get(Sunday_Reading_URL)
    else:
        browser.get(Today_Reading_URL)

    timeout = 60000

    try:
        WebDriverWait(browser, timeout).until(EC.visibility_of_element_located((By.XPATH, "//p")))
        print(browser.page_source)
        print("This ran")
    except TimeoutException:
        print("Timed out waiting for page to load")
        browser.quit()
    readings = browser.find_element_by_id("ci-readings")
    final = readings.text.replace(u'\u2019', '\'').replace(u'\u2018', '\'').replace(u'\u201c', '\"').replace(u'\u201d', '\"')
    browser.close()


    # Wait 30 seconds for page to load
    #timeout = 30
    #try:
    #    #WebDriverWait(browser, timeout).until(EC.visibility_of_element_located((By.XPATH, "//span[@style=\"font-family: Times New Roman;\"]")))
    #    WebDriverWait(browser, timeout).until(EC.visibility_of_element_located((By.XPATH, "//p")))
    #except TimeoutException:
    #    print("Timed out waiting for page to load")
    #    browser.quit()



    splitstring = final.split()
    readings = []
    gospel = []
    response = []
    count = 0
    gospel_count = 0
    reading_count = 0
    response_count = 0
    read_start_count = read_end_count = 0
    gosp_start_count = gosp_end_count = 0
    resp_start_count = resp_end_count  = 0
    cur_reading = False
    gosp_reading = False
    resp_reading = False
    for word in splitstring:

        if (cur_reading==False and word.lower()=="reading"):
            cur_reading = True
            read_start_count = count+1

        if (resp_reading==False and word.lower()=="responsorial"):
            resp_reading = True
            resp_start_count = count+2

        if (gosp_reading==False and word.lower()=="gospel"):
            if(splitstring[count+1].lower() == "acclamation"):
                count = count + 1
                continue
            elif(splitstring[count+1].lower() == "reflection"):
                break
            gosp_reading = True
            gosp_start_count = count+1


        if(cur_reading==True and splitstring[count].lower()=="the" and splitstring[count+1].lower()=="word" and splitstring[count+2].lower()=="of" and splitstring[count+3].lower()=="the" and (splitstring[count+4].lower()=="lord" or splitstring[count+4].lower()=="lord.")):
            readings.append("")
            cur_reading = False
            read_end_count = count
            for i in range(read_start_count, read_end_count):
                readings[reading_count] = readings[reading_count] + splitstring[i] + " "
            readings[reading_count] = readings[reading_count] + "\n\n" + "The Word of the Lord."
            reading_count = reading_count + 1

        if(gosp_reading==True and splitstring[count].lower()=="the" and splitstring[count+1].lower()=="gospel" and splitstring[count+2].lower()=="of" and splitstring[count+3].lower()=="the" and (splitstring[count+4].lower()=="lord" or splitstring[count+4].lower()=="lord.")):
            gospel.append("")
            gosp_reading = False
            gosp_end_count = count
            for i in range(gosp_start_count, gosp_end_count):
                gospel[gospel_count] = gospel[gospel_count] + splitstring[i] + " "
            gospel[gospel_count] = gospel[gospel_count] +"\n\n" + "The Gospel of the Lord." 
            gospel_count = gospel_count + 1

        if(resp_reading==True and (splitstring[count].lower()=="second" or word.lower()=="gospel" or splitstring[count].lower()=="first" or splitstring[count].lower()=="third" or splitstring[count].lower()=="fourth" or splitstring[count].lower()=="fifth")):
            response.append("")
            resp_reading = False
            resp_end_count = count
            for i in range(resp_start_count, resp_end_count):
                if(splitstring[i][0].isdigit()):
                    response[response_count] = response[response_count] + "\n\n"
                response[response_count] = response[response_count] + splitstring[i] + " "

            response_count = response_count + 1

        count = count + 1

    


    return jsonify(readings=readings, gospels=gospel, response=response)

"""def duplicate_slide(prayer_slide_no, prs):
        template = Prayers_Pres.slides[prayer_slide_no]
        try:
            blank_slide_layout = Prayers_Pres.slide_layouts[prayer_slide_no]
        except:
            print("Do nothing")

        copied_slide = prs[0].slides.add_slide(blank_slide_layout)

        #for shp in template.shapes:
        #    el = shp.element
        #    newel = copy.deepcopy(el)
        #    copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        #for _, value in six.iteritems(template.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
        #    if "notesSlide" not in value.reltype:
        #        copied_slide.part.rels.add_relationship(value.reltype, value._target, value.rId)

        return copied_slide """

"""def Generate_Prayer_Slide(BackgroundPicturePath, TitleFontColour, ContentFontColour, title_place_holder, content_place_holder, slide_no, prs):
    temp_slide = Prayers_Pres.slides[slide_no]
    for shape in temp_slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
    new_slide_layout = prs[0].slide_layouts[1]
    slide = prs[0].slides.add_slide(new_slide_layout)

    #slide = prs[0].slides.add_slide(duplicate_slide(slide_no, prs))

    new_slide_title = slide.shapes.title
    #new_slide_content = slide.placeholders[1]

    Background_Image(slide, BackgroundPicturePath, prs)
    if title_place_holder != "":
        temp_title = temp_slide.placeholders[title_place_holder]
    temp_content = temp_slide.placeholders[content_place_holder]

    if title_place_holder != "":
        new_slide_content = slide.shapes.add_textbox(0, 0, prs[0].slide_width, prs[0].slide_height)
    else:
        new_slide_content = slide.shapes.add_textbox(0, Inches(0.8), prs[0].slide_width, prs[0].slide_height - Inches(0.8))
    

    try:
        temp_font_title = temp_title.text_frame
        new_slide_title.text_frame = temp_font_title
        new_slide_tile.textframe.paragraph[0].font.color.rgb = RGBColor(int(TitleFontColour[0]), int(TitleFontColour[1]), int(TitleFontColour[2]))

    except Exception:
        print("Do nothing")

    temp_font_content = temp_content.text_frame
    new_slide_content.text_frame = temp_font_content
    new_slide_content.textframe.paragraph[0].font.color.rgb = RGBColor(int(ContentFontColour[0]), int(ContentFontColour[1]), int(ContentFontColour[2]))"""



def Generate_Unit_Slide(title_content, content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, alignment):
    
    slide_layout = prs[0].slide_layouts[1]
    
    slide = prs[0].slides.add_slide(slide_layout)
    #body = slide.placeholders[1]
    if title_content=="":
        body = slide.shapes.add_textbox(0, 0, prs[0].slide_width, prs[0].slide_height)
    else:
        body = slide.shapes.add_textbox(0, Inches(2), prs[0].slide_width, prs[0].slide_height - Inches(2))
    Background_Image(slide, BackgroundPicturePath, prs)
    try:
        title = slide.shapes.title
    except Exception:
        print("Title Holder does not Exist")
    
    try:
        tf = title.text_frame
        tf.word_wrap = True
        #tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title_content
        font = run.font
        font.name = 'Calibri'
        font.bold = True
        font.size = Pt(60)
        font.color.rgb = RGBColor(int(TitleFontColour[0]), int(TitleFontColour[1]), int(TitleFontColour[2]))
    except Exception:
        print("Title Holder does not exist")
    
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    #tf.text = content
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.text = content
    font = p.font
    font.name = 'Calibri'
    font.bold = True
    font.size = Pt(50)
    p.font.color.rgb = RGBColor(int(ContentFontColour[0]), int(ContentFontColour[1]), int(ContentFontColour[2]))
    
    


def Generate_Slides_For_Large_Data(large_data, title, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs):

    
    large_data = large_data.strip()
    n_large_data = large_data.split('\n\n')
    paras = []

    for i in n_large_data:
        temp = i.split('\r\n')
        for j in temp:
            paras.append(j)
    
    
    title_data = title
    slide_content = ""
    abs_count = 0
    for para in paras:
        """temp_lines = para.splitlines()
        for s in temp_lines:
            lines = s.split('.')"""
        lines = para.splitlines()
        print (lines)
        count = 0
        count_check = 4
        for line in lines:
            if line == '':
                continue
            slide_content = slide_content + line + "\n"
            if abs_count == 0:
                count_check = 2
            else:
                count_check = 4
            if(count == count_check):
                #print (slide_content)
                Generate_Unit_Slide(title_data, slide_content.strip(), BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)
                title_data = ""
                slide_content = ""
                count = 0
                abs_count = abs_count + 1
            else:
                count = count + 1
        if(count != 0):
            #print(slide_content)
            Generate_Unit_Slide(title_data, slide_content.strip(), BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)
            slide_content = ""
            title_data = ""
    """for i in range(0, len(paras)):
        if i == 0:
            title_data = title
        else:
            title_data = ""
        words_in_para = paras[i].split()
        while True:
            if len(words_in_para) > 30 and len(words_in_para[31:]) < 40:
                for word in words_in_para:
                    slide_content = slide_content + word + " "
                words_in_para = words_in_para[31:]
                Generate_Unit_Slide(title_data, slide_content.strip(), BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)
            else:
                for word in words_in_para:
                    slide_content = slide_content + word + " "
                Generate_Unit_Slide(title_data, slide_content.strip(), BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)
                break"""
        
def List_Splitter(list_of_songs, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs):
    for data in list_of_songs:
        print(data['title'])
        Generate_Slides_For_Large_Data(data['content'], data['title'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)
        Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)
    




def Generate_PPT(ppt_datax):
    print(ppt_datax)
    ppt_data = json.loads(json.dumps(ppt_datax))
    prs = [Presentation()]
    #Global Parameters
    text_colours = []
    BackgroundPicturePath = "Black_Background_Default.jpg"
    BackgroundPicturePath = os.path.join(path_to_here, BackgroundPicturePath)
    if ppt_data['does_background_exist'] == True:
        fbst.child(ppt_data['background_image_path']).download(os.path.join(path_to_here, ppt_data['background_image_path']))
        BackgroundPicturePath = os.path.join(path_to_here, ppt_data['background_image_path'])
        bgimage = cv2.imread(BackgroundPicturePath)
        width = 960
        height = 720
        dim = (width, height)
        resized = cv2.resize(bgimage, dim, interpolation = cv2.INTER_CUBIC)
        cv2.imwrite(BackgroundPicturePath, resized)
    TitleFontColour = ContentFontColour = []
    if ppt_data['use_generated_colours'] == True:
        text_colours = get_dominant_color(cv2.cvtColor(cv2.imread(BackgroundPicturePath), cv2.COLOR_BGR2HSV))
        text_colours = complement(int(text_colours[0]), int(text_colours[1]), int(text_colours[2]))
        TitleFontColour = text_colours
        ContentFontColour = text_colours
    else:
        TitleFontColour = ppt_data['title_colour']
        ContentFontColour = ppt_data['content_colour']

    #Title Page
    MainTitlePageContent = ppt_data["main_title"]
    MainTitlePageSubtitle = ppt_data["main_subtitle"]


    """title_slide_layout = prs[0].slide_layouts[0]
    slide = prs[0].slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    Background_Image(slide, BackgroundPicturePath, prs)
    tf = title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = MainTitlePageContent
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(80)
    font.bold = True
    font.color.rgb = RGBColor(int(TitleFontColour[0]), int(TitleFontColour[1]), int(TitleFontColour[2]))


    tf = subtitle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.add_paragraph()

    run = p.add_run()
    run.text = MainTitlePageSubtitle
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(80)
    font.bold = True
    run.font.color.rgb = RGBColor(int(ContentFontColour[0]), int(ContentFontColour[1]), int(ContentFontColour[2]))"""

    Generate_Unit_Slide(MainTitlePageContent, MainTitlePageSubtitle, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    List_Splitter(ppt_data['entrance_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    #Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    #Generate_Prayer_Slide(BackgroundPicturePath, TitleFontColour, ContentFontColour, 0, 1, 0, prs)
    
    prayer_title = "The Introductory Rite"
    prayer_content = "Priest:\nIn the name of the Father, and of the Son, and of the Holy Spirit.\n\nPeople:\nAmen."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_content = "Priest:\nThe grace of our Lord Jesus Christ, and the love of God, and the communion of the Holy Spirit be with you all.\n\nPeople:\nAnd with your spirit. "
    prayer_title = ""

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)
    
    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    pen_title = "The Penitential Act"
    pen_content = ""
    prayer_content = "Priest:\nDear Brothers and Sisters let us acknowledge our sins, and so prepare ourselves to celebrate the sacred mysteries.\n\n"
    
    pen_content = pen_content + prayer_content

    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    #prayer_title = ""
    prayer_content = "I confess to almighty God\nand to you, my brothers and sisters,\nthat I have greatly sinned, in my thoughts and\nin my words, in what I have done\nand in what I have failed to do,\n(And, striking their breast, they say)\nthrough my fault, through my\nfault, through my most grievous fault;"
    pen_content = pen_content + prayer_content
    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    #prayer_title = ""
    prayer_content = "therefore I ask blessed Mary\never-Virgin, all the Angels and Saints, and you,\nmy brothers and sisters,\nto pray for me to the Lord our God.\n\nThe absolution by the Priest follows:\nMay almighty God have mercy on us, forgive us our sins, and bring us to everlasting life. Amen."
    pen_content = pen_content + prayer_content

    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    Generate_Slides_For_Large_Data(pen_content, pen_title, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    List_Splitter(ppt_data['kyrie_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    List_Splitter(ppt_data['gloria_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    List_Splitter(ppt_data['reading_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    List_Splitter(ppt_data['acclamation_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    prayer_title = ""
    prayer_content = "Priest:\nThe Lord be with you.\n\nPeople:\nAnd with your spirit."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nA reading from the Holy Gospel  according to St. ________\n\nPeople:\nGlory to you, O Lord."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)


    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    List_Splitter(ppt_data['gospel_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    prayer_title = "The Apostle's Creed"
    prayer_content = "I believe in God, the Father\nAlmighty, Creator of heaven and\nearth, and in Jesus Christ, His only\nSon, our Lord, who was "

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "conceived by the Holy Spirit,\nborn of the Virgin Mary,\nsuffered under Pontius Pilate,\nwas crucified, died and was buried;\nHe descended into hell;\non the third day He rose again from"

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "the dead; He ascended into\nheaven, and is seated at the right \nhand of God the Father Almighty; \nfrom there He will come to judge \nthe living and the dead."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "I believe in the Holy Spirit,\nthe Holy Catholic Church,\nthe communion of Saints,\nthe forgiveness of sins,\nthe resurrection of the body,\nand life everlasting. Amen."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    nic_title = "The Nicene Creed"
    #prayer_title = "The Nicene Creed"
    nic_content = ""
    prayer_content = "I believe in one God, the Father\nalmighty, maker of heaven and earth,\nof all things visible and invisible. I\nbelieve in one Lord Jesus Christ, the \nOnly Begotten Son of God,"
    nic_content = nic_content + prayer_content

    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    #prayer_title = ""
    prayer_content = " born of the Father\nbefore all ages.\nGod from God, Light from Light, true\nGod from true God, begotten, not\nmade, consubstantial with the Father;\nthrough him all things were made."
    nic_content = nic_content + prayer_content
    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    #prayer_title = ""
    prayer_content = "For us men and\nfor our salvation he\ncame down from heaven,\nand by the Holy\nSpirit was incarnate\nof the Virgin Mary,\nand became man.\nFor our sake he\nwas crucified under\nPontius Pilate, "
    nic_content = nic_content + prayer_content
    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    #prayer_title = ""
    prayer_content = "he suffered death and was buried,\nand rose again on\n the third day in\naccordance with the Scriptures.\nHe ascended into heaven and is\nseated at the right hand of the\nFather. He will come again in glory"
    nic_content = nic_content + prayer_content
    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    #prayer_title = ""
    prayer_content = "to judge the living\nand the dead\nand his kingdom will\nhave no end.\nI believe in the Holy Spirit, the Lord,\nthe giver of life,\nwho proceeds from\nthe Father and the Son,\nwho with the Father and the Son is"
    nic_content = nic_content + prayer_content

    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    #prayer_title = ""
    prayer_content = "adored and glorified,\nwho has spoken\nthrough the prophets.\nI believe in one,\nholy, catholic and apostolic Church.\nI confess one baptism\nfor the forgiveness of\nsins and I look forward\nto the resurrection\nof the dead and the life of the world to\ncome. Amen!"
    nic_content = nic_content + prayer_content

    #Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    Generate_Slides_For_Large_Data(nic_content, nic_title, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)


    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)


    prayer_title = "Prayer of the Faithful"
    prayer_content = "Your Response:\n"

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    List_Splitter(ppt_data['offertory_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    prayer_title = ""
    prayer_content = "Priest:\nPray brothers and sisters, that my sacrifice and yours may be acceptable to God, the almighty Father."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "People:\nMay the Lord accept the sacrifice at your hands for the praise and glory of his name, for our good and the good of all his holy Church."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = "The Eucharistic Prayer"
    prayer_content = "Priest:\nThe Lord be with you.\n\nPeople:\nAnd with your spirit."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nLift up your hearts.\n\nPeople:\nWe lift them up to the Lord."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nLet us give thanks to the Lord our God.\n\nPeople:\nIt is right and just."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    
    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    List_Splitter(ppt_data['holy_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    List_Splitter(ppt_data['proclamation_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    prayer_title = ""
    prayer_content = "Priest:\nThrough him, and with him, and in him, O God, almighty Father, in the unity of the Holy Spirit, all glory and honour is yours, for ever and ever.\n\nPeople:\nAmen."

    Generate_Slides_For_Large_Data(prayer_content, prayer_title, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nAt the Saviour's command and formed by divine teaching, we dare to say:"

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = "The Lord's Prayer"
    prayer_content = "Our Father, who art in Heaven, hallowed be thy name; thy kingdom come, thy will be done on earth as it is in Heaven.\n\nGive us this day our daily bread,\nand forgive us our trespasses,\nas we forgive those who trespass against us;\nand lead us not into temptation,\nbut deliver us from evil."
    Generate_Slides_For_Large_Data(prayer_content, prayer_title, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    prayer_title = ""
    prayer_content = "Priest:\nDeliver us, Lord, we pray, from every evil,\ngraciously grant peace in our days,\nthat, by the help of your mercy,\nwe may be always free\nfrom sin and safe from all\ndistress, as we await the blessed\nhope and the coming of our Saviour, Jesus Christ."

    Generate_Slides_For_Large_Data(prayer_content, prayer_title, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    prayer_title = ""
    prayer_content = "The Priest and the people conclude the prayer, acclaiming:\nFor the kingdom, the power and the glory are yours now and for ever."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nLord Jesus Christ, who said to your Apostles:\nPeace I leave you, my peace I give you;\nlook not on our sins,\nbut on the faith of your Church, and graciously grant her peace and unity in accordance with your will, who live and reign for ever and ever.\n\nPeople:\nAmen."

    Generate_Slides_For_Large_Data(prayer_content, prayer_title, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)

    prayer_title = ""
    prayer_content = "Priest:\nThe peace of the Lord be with you always.\n\nPeople:\nAnd with your spirit."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nLet us offer each other the sign of peace. Peace be with you."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)


    prayer_title = ""
    prayer_content = "Lamb of God You take away\nThe sins of the world,\nhave mercy on us (2)\nLamb of God You take away\nThe sins of the world\nGrant us peace"

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nBehold the Lamb of God, Behold him who takes away the sins of the world, Blessed are those called to the Supper of the Lamb."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest and People:\nLord I am not worthy that  you should enter under my roof, But only say the word and my soul shall be healed."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    
    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)


    List_Splitter(ppt_data['communion_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)


    prayer_title = ""
    prayer_content = "Announcements for the Week"

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    
    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = "The Concluding Rites"
    prayer_content = "Priest:\nThe Lord be with you.\n\nPeople:\nAnd with your spirit."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nMay almighty God bless you, the Father, and the Son, and the Holy Spirit.\n\nPeople:\nAmen."

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    prayer_title = ""
    prayer_content = "Priest:\nGo and announce the Gospel of the Lord.\n\nPeople:\nThanks be to God"

    Generate_Unit_Slide(prayer_title, prayer_content, BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)

    Generate_Unit_Slide("", "", BackgroundPicturePath, TitleFontColour, ContentFontColour, prs, PP_ALIGN.CENTER)


    List_Splitter(ppt_data['recessional_list'], BackgroundPicturePath, TitleFontColour, ContentFontColour, prs)


    #Saving the ppt
    file_path = str(uuid.uuid1()) + ".pptx" #(ppt_data["background_image_path"][7:ppt_data["background_image_path"].rindex('.')] + ".pptx").replace(':','_')
    print(file_path)
    prs[0].save(os.path.join(path_to_here, file_path))

    fbst.child("files/"+file_path).put(os.path.join(path_to_here, file_path), user['idToken'])

    os.remove(os.path.join(path_to_here, file_path))

    return ("files/"+file_path)



app = Flask(__name__)



@app.route("/")
def home():
    return "You've reached the server. Dance around now!"

@app.route("/lyrics")
def lyrics():
    lyquery = flrequest.args.get('lyquery')
    return jsonify(lyrics=Lyric_Resultify(lyquery))

@app.route("/readings")
def readings():
    rdchoice = flrequest.args.get('rdchoice')
    return Reading_Resultify(rdchoice)

@app.route("/createppt", methods=['GET', 'POST'])
def createppt():
    if(flrequest.method == 'POST'):
        ppt_data = flrequest.get_json()
        return jsonify(download_url=Generate_PPT(ppt_data))

if __name__ == "__main__":
    app.run()
